# backend/utils/parsers.py

import os
from pdfplumber import open as open_pdf
from pptx import Presentation
import pandas as pd
from docx import Document

def parse_pdf(path):
    pages = []
    with open_pdf(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            pages.append({
                "text": text,
                "source": {"type": "pdf", "page": i}
            })
    return pages

def parse_pptx(path):
    slides = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        slides.append({
            "text": text,
            "source": {"type": "pptx", "page": i}   # unify on "page"
        })
    return slides

def parse_csv(path):
    df = pd.read_csv(path)
    outputs = []
    for i, row in df.iterrows():
        text = row.to_csv(header=False).strip()
        outputs.append({
            "text": text,
            "source": {"type": "csv", "row": int(i)}
        })
    return outputs

def parse_docx(path):
    doc = Document(path)
    paras = []
    for i, para in enumerate(doc.paragraphs, start=1):
        text = para.text
        if not text:
            continue
        paras.append({
            "text": text,
            "source": {"type": "docx", "paragraph": i}
        })
    return paras

def parse_txt(path):
    outputs = []
    with open(path, encoding="utf-8") as f:
        for i, line in enumerate(f, start=1):
            text = line.strip()
            if text:
                outputs.append({
                    "text": text,
                    "source": {"type": "txt", "line": i}
                })
    return outputs

# Dispatcher
def parse_file(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return parse_pdf(path)
    if ext in [".pptx", ".ppt"]:
        return parse_pptx(path)
    if ext == ".csv":
        return parse_csv(path)
    if ext in [".docx", ".doc"]:
        return parse_docx(path)
    if ext in [".txt", ".md"]:
        return parse_txt(path)
    raise ValueError(f"Unsupported file type: {ext}")
